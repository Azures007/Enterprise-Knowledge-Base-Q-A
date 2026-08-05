"""
=============================================================================
文档加载模块

支持多种文档格式的解析与文本提取：

    📄 PDF (.pdf)
        - 正常 PDF → PyMuPDF 按页提取
        - 扫描版 PDF → 自动 OCR 兜底（需 PaddleOCR）
        - 加密 PDF → 检测加密并提示
        - 复杂排版 → 按坐标块重排

    📝 Office 文档
        - Word (.docx)  → python-docx + XML 增强
        - Word (.doc)   → textract / olefile / LibreOffice 转换
        - PowerPoint (.pptx) → python-pptx
        - Excel (.xlsx)  → openpyxl
        - Excel (.xls)   → xlrd

    🖼️ 图片 (.jpg/.jpeg/.png/.bmp)
        - OCR 文字识别（PaddleOCR / EasyOCR）

    📄 纯文本类
        - TXT, MD, PY, JSON, YAML, CSV, HTML, XML
        - 自动检测编码（chardet）+ 多编码备选

    📦 WPS 格式 (.wps/.et)
        - 通过 LibreOffice 命令行转换为 docx/xlsx 后解析

    ⚠️ 超大文档检测与分批处理
=============================================================================
"""

import json
import os
import re
import shutil
import subprocess
import tempfile
import uuid
from pathlib import Path
from typing import Any, Optional

from src.utils.logger import setup_logger

logger = setup_logger(__name__)

LARGE_FILE_THRESHOLD = 20 * 1024 * 1024  # 20MB
HUGE_FILE_THRESHOLD = 100 * 1024 * 1024  # 100MB


class DocumentLoaderError(Exception):
    pass


class DocumentLoader:
    SUPPORTED_EXTENSIONS = {
        ".pdf": "_load_pdf",
        ".docx": "_load_docx",
        ".doc": "_load_doc",
        ".docm": "_load_docx",
        ".pptx": "_load_pptx",
        ".xlsx": "_load_excel",
        ".xls": "_load_excel",
        ".jpg": "_load_image",
        ".jpeg": "_load_image",
        ".png": "_load_image",
        ".bmp": "_load_image",
        ".tiff": "_load_image",
        ".tif": "_load_image",
        ".txt": "_load_text",
        ".md": "_load_text",
        ".py": "_load_text",
        ".yaml": "_load_text",
        ".yml": "_load_text",
        ".html": "_load_text",
        ".htm": "_load_text",
        ".xml": "_load_text",
        ".csv": "_load_text",
        ".json": "_load_json",
        ".wps": "_load_wps",
        ".et": "_load_wps",
    }

    def __init__(self, default_password=None):
        self._load_count = 0
        self._error_count = 0
        self.default_password = default_password

    def load_file(self, file_path, password=None):
        path = Path(file_path)
        if not path.exists():
            raise DocumentLoaderError(f"文件不存在: {path}")
        file_size = path.stat().st_size
        if file_size > HUGE_FILE_THRESHOLD:
            raise DocumentLoaderError(f"文件过大 ({file_size/1024/1024:.1f}MB)")
        if file_size > LARGE_FILE_THRESHOLD:
            logger.warning(f"文件较大 ({file_size/1024/1024:.1f}MB)")
        ext = path.suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise DocumentLoaderError(f"不支持的文件格式 '{ext}'")
        method_name = self.SUPPORTED_EXTENSIONS[ext]
        method = getattr(self, method_name)
        try:
            logger.info(f"正在解析文档: {path.name} (格式: {ext})")
            documents = method(path, password)
            self._load_count += 1
            logger.info(f"文档解析完成: {path.name} → {len(documents)} 个文档块")
            return documents
        except DocumentLoaderError:
            raise
        except Exception as e:
            self._error_count += 1
            raise DocumentLoaderError(f"文档解析失败 [{path.name}]: {e}") from e

    def load_directory(self, directory, recursive=True):
        dir_path = Path(directory)
        if not dir_path.is_dir():
            raise DocumentLoaderError(f"目录不存在: {directory}")
        pattern = "**/*" if recursive else "*"
        all_docs = []
        for file_path in sorted(dir_path.glob(pattern)):
            if not file_path.is_file():
                continue
            ext = file_path.suffix.lower()
            if ext not in self.SUPPORTED_EXTENSIONS:
                continue
            try:
                docs = self.load_file(file_path)
                all_docs.extend(docs)
            except DocumentLoaderError as e:
                logger.warning(str(e))
                continue
        logger.info(f"批量加载完成: 成功 {self._load_count} 个文件, 失败 {self._error_count} 个文件, 共 {len(all_docs)} 个文档块")
        return all_docs

    @staticmethod
    def _load_pdf(path, password=None):
        import fitz
        documents = []
        try:
            doc = fitz.open(str(path))
        except Exception as e:
            raise DocumentLoaderError(f"无法打开 PDF: {e}")
        if doc.is_encrypted:
            if password:
                if not doc.authenticate(password):
                    doc.close()
                    raise DocumentLoaderError("PDF 密码错误")
            else:
                doc.close()
                raise DocumentLoaderError("PDF 已加密，请提供密码")
        total_pages = len(doc)
        for page_num in range(total_pages):
            page = doc[page_num]
            text = page.get_text().strip()
            if len(text) < 20:
                blocks = page.get_text("blocks")
                if blocks:
                    blocks.sort(key=lambda b: (round(b[1], -1), b[0]))
                    lines = [b[4].strip() for b in blocks if len(b) > 4 and b[4].strip()]
                    text = "\n".join(lines)
            if len(text) < 10:
                try:
                    pix = page.get_pixmap(dpi=300)
                    from src.document_loader.ocr import ocr_pdf_page
                    ocr_text = ocr_pdf_page(pix)
                    if ocr_text:
                        text = ocr_text
                except Exception as e:
                    logger.warning(f"第 {page_num+1} 页 OCR 失败: {e}")
                    text = ""
            if not text:
                continue
            documents.append({
                "page_content": text,
                "metadata": {
                    "source": str(path), "filename": path.name,
                    "page": page_num + 1, "total_pages": total_pages,
                    "file_type": "pdf", "file_size": path.stat().st_size,
                },
            })
        doc.close()
        return documents

    # ================================================================
    # .docx 解析（拆分为多个独立方法，提升可读性与可测试性）
    # ================================================================

    @staticmethod
    def _load_docx(path, password=None):
        """解析 .docx 文档，提取正文、页眉页脚、脚注尾注、批注、表格等内容。"""
        from docx import Document as DocxDocument

        try:
            doc = DocxDocument(str(path))
        except Exception as e:
            raise DocumentLoaderError(f"无法打开 docx 文件: {e}")

        raw_xmls = DocumentLoader._read_docx_xmls(path)
        warnings_list = []

        all_text_parts = []
        paragraphs_text = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        if paragraphs_text:
            all_text_parts.append("[正文]")
            all_text_parts.append("\n".join(paragraphs_text))

        all_text_parts.extend(DocumentLoader._extract_docx_headers_footers(doc))
        all_text_parts.extend(DocumentLoader._extract_docx_notes(raw_xmls))
        all_text_parts.extend(DocumentLoader._extract_docx_comments(raw_xmls))
        all_text_parts.extend(DocumentLoader._extract_docx_textboxes(raw_xmls))
        all_text_parts.extend(DocumentLoader._extract_docx_revisions(raw_xmls))
        all_text_parts.extend(DocumentLoader._extract_docx_tables(raw_xmls, doc))

        # 检查嵌入对象
        try:
            import zipfile
            with zipfile.ZipFile(str(path), "r") as z:
                embedded = [n for n in z.namelist() if "embeddings" in n]
                if embedded:
                    warnings_list.append(
                        f"文档包含 {len(embedded)} 个嵌入文件，嵌入内容无法提取，请单独上传"
                    )
        except Exception:
            pass

        full_text = "\n\n".join(all_text_parts) if all_text_parts else ""
        if not full_text.strip():
            warnings_list.append("文档中未提取到文本内容")

        return [{
            "page_content": full_text,
            "metadata": {
                "source": str(path), "filename": path.name,
                "file_type": path.suffix.lower().lstrip("."),
                "paragraphs": len(paragraphs_text),
                "file_size": path.stat().st_size,
                "warnings": warnings_list if warnings_list else None,
            },
        }]

    # ----------------------------------------------------------
    # .docx 辅助方法
    # ----------------------------------------------------------

    @staticmethod
    def _read_docx_xmls(path) -> dict:
        """读取 docx 包内所有 .xml 文件的原始字节，供深度解析使用。"""
        import zipfile
        raw_xmls = {}
        try:
            with zipfile.ZipFile(str(path), "r") as z:
                for name in z.namelist():
                    if name.endswith(".xml"):
                        raw_xmls[name] = z.read(name)
        except Exception:
            pass
        return raw_xmls

    @staticmethod
    def _docx_ns(tag: str) -> str:
        """将 WordprocessingML 简写标签转换为完整命名空间标签。"""
        return f"{{http://schemas.openxmlformats.org/wordprocessingml/2006/main}}{tag}"

    @staticmethod
    def _docx_extract_text(elem) -> str:
        """从 XML 元素中提取全部 w:t 文本。"""
        texts = []
        for t in elem.iter(DocumentLoader._docx_ns("t")):
            if t.text:
                texts.append(t.text)
        return "".join(texts)

    @staticmethod
    def _extract_docx_headers_footers(doc) -> list[str]:
        """提取页眉页脚文本。"""
        parts = []
        header_texts = []
        footer_texts = []
        try:
            for section in doc.sections:
                for h in (section.header, section.footer):
                    if not h:
                        continue
                    for p in h.paragraphs:
                        t = p.text.strip()
                        if t:
                            (header_texts if h is section.header else footer_texts).append(t)
        except Exception:
            pass
        if header_texts:
            parts.append("[页眉]")
            parts.append("\n".join(header_texts))
        if footer_texts:
            parts.append("[页脚]")
            parts.append("\n".join(footer_texts))
        return parts

    @staticmethod
    def _extract_docx_notes(raw_xmls: dict) -> list[str]:
        """提取脚注（footnotes）与尾注（endnotes）。"""
        parts = []
        ET = DocumentLoader._get_etree()
        ns = DocumentLoader._docx_ns

        for key, label in (("word/footnotes.xml", "脚注"), ("word/endnotes.xml", "尾注")):
            xml_data = raw_xmls.get(key, b"")
            if not xml_data:
                continue
            try:
                root = ET.fromstring(xml_data)
                tag = "footnote" if "footnote" in key else "endnote"
                texts = [
                    DocumentLoader._docx_extract_text(e).strip()
                    for e in root.findall(f".//{ns(tag)}")
                ]
                texts = [t for t in texts if t]
                if texts:
                    parts.append(f"[{label}]")
                    parts.append("\n".join(texts))
            except Exception:
                continue
        return parts

    @staticmethod
    def _extract_docx_comments(raw_xmls: dict) -> list[str]:
        """提取批注（comments）。"""
        parts = []
        ET = DocumentLoader._get_etree()
        ns = DocumentLoader._docx_ns
        xml_data = raw_xmls.get("word/comments.xml", b"")
        if not xml_data:
            return parts
        try:
            root = ET.fromstring(xml_data)
            cm_texts = []
            for cm in root.findall(f".//{ns('comment')}"):
                author = cm.get("author", "未知")
                txt = DocumentLoader._docx_extract_text(cm)
                if txt.strip():
                    cm_texts.append(f"[批注-{author}]: {txt.strip()}")
            if cm_texts:
                parts.append("[批注]")
                parts.append("\n".join(cm_texts))
        except Exception:
            pass
        return parts

    @staticmethod
    def _extract_docx_textboxes(raw_xmls: dict) -> list[str]:
        """提取浮动文本框（txbxContent）。"""
        parts = []
        ET = DocumentLoader._get_etree()
        ns = DocumentLoader._docx_ns
        xml_data = raw_xmls.get("word/document.xml", b"")
        if not xml_data:
            return parts
        try:
            root = ET.fromstring(xml_data)
            txbx_texts = [
                DocumentLoader._docx_extract_text(txbx).strip()
                for txbx in root.iter(ns("txbxContent"))
            ]
            txbx_texts = [t for t in txbx_texts if t]
            if txbx_texts:
                parts.append("[浮动文本框]")
                parts.append("\n".join(txbx_texts))
        except Exception:
            pass
        return parts

    @staticmethod
    def _extract_docx_revisions(raw_xmls: dict) -> list[str]:
        """提取修订插入内容（ins）。"""
        parts = []
        ET = DocumentLoader._get_etree()
        ns = DocumentLoader._docx_ns
        xml_data = raw_xmls.get("word/document.xml", b"")
        if not xml_data:
            return parts
        try:
            root = ET.fromstring(xml_data)
            ins_texts = []
            for ins in root.iter(ns("ins")):
                author = ins.get("author", "未知")
                txt = DocumentLoader._docx_extract_text(ins)
                if txt.strip():
                    ins_texts.append(f"[修订插入-{author}]: {txt.strip()}")
            if ins_texts:
                parts.append("[修订内容]")
                parts.append("\n".join(ins_texts))
        except Exception:
            pass
        return parts

    @staticmethod
    def _extract_docx_tables(raw_xmls: dict, doc) -> list[str]:
        """
        提取表格内容。

        优先用原始 XML（保留合并单元格语义），失败时回退到 python-docx 的 tables。
        """
        parts = []
        ET = DocumentLoader._get_etree()
        ns = DocumentLoader._docx_ns
        xml_data = raw_xmls.get("word/document.xml", b"")
        if xml_data:
            try:
                root = ET.fromstring(xml_data)
                table_idx = 0
                for tbl in root.iter(ns("tbl")):
                    table_idx += 1
                    table_rows = []
                    for tr in tbl.iter(ns("tr")):
                        cells = []
                        for tc in tr.iter(ns("tc")):
                            tc_pr = tc.find(ns("tcPr"))
                            span = 1
                            if tc_pr is not None:
                                gs = tc_pr.find(ns("gridSpan"))
                                if gs is not None and gs.get("val"):
                                    span = int(gs.get("val"))
                                vm = tc_pr.find(ns("vMerge"))
                                if vm is not None and vm.get("val") != "restart":
                                    continue
                            cell_text = DocumentLoader._docx_extract_text(tc)
                            cells.append(cell_text.strip())
                        if any(c for c in cells):
                            table_rows.append(" | ".join(cells))
                    if table_rows:
                        parts.append(f"[表格 {table_idx}]")
                        parts.append("\n".join(table_rows))
                if parts:
                    return parts
            except Exception:
                parts = []

        # 回退：python-docx 的 tables
        try:
            tables_text = []
            for table in doc.tables:
                rows_text = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
                tables_text.append("\n".join(rows_text))
            if tables_text:
                parts.append("[表格内容]")
                parts.append("\n\n".join(tables_text))
        except Exception:
            pass
        return parts

    @staticmethod
    def _get_etree():
        """延迟导入 xml.etree.ElementTree"""
        from xml.etree import ElementTree as ET
        return ET

    @staticmethod
    def _try_extract_doc_text(path):
        try:
            import textract
            text = textract.process(str(path)).decode("utf-8", errors="ignore")
            if text and text.strip():
                logger.info(f"使用 textract 解析 .doc 成功: {path.name}")
                return text.strip()
        except ImportError:
            pass
        except Exception as e:
            logger.warning(f"textract 解析失败: {e}")
        try:
            import olefile
            if olefile.isOleFile(str(path)):
                ole = olefile.OleFileIO(str(path))
                try:
                    if ole.exists("WordDocument"):
                        data = ole.openstream("WordDocument").read()
                        text = data.decode("utf-16-le", errors="ignore")
                        text = re.sub(r"[^一-鿿 -~\n\r]", "", text)
                        text = re.sub(r"\s+", " ", text).strip()
                        if text and len(text) > 20:
                            logger.info(f"使用 olefile 解析 .doc 成功: {path.name}")
                            return text
                finally:
                    ole.close()
        except ImportError:
            pass
        except Exception as e:
            logger.warning(f"olefile 解析失败: {e}")
        return None

    @staticmethod
    def _find_libreoffice():
        for cmd in ["libreoffice", "soffice"]:
            try:
                subprocess.run([cmd, "--version"], capture_output=True, timeout=10)
                return cmd
            except (FileNotFoundError, subprocess.TimeoutExpired):
                pass
        common_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            os.path.expanduser(r"~\AppData\Local\Programs\LibreOffice\program\soffice.exe"),
        ]
        for p in common_paths:
            if os.path.exists(p):
                return p
        return None

    @staticmethod
    def _load_doc(path, password=None):
        doc_text = DocumentLoader._try_extract_doc_text(path)
        if doc_text:
            return [{
                "page_content": doc_text,
                "metadata": {
                    "source": str(path), "filename": path.name,
                    "file_type": "doc", "file_size": path.stat().st_size,
                },
            }]
        lo_found = DocumentLoader._find_libreoffice()
        if not lo_found:
            raise DocumentLoaderError(
                "旧版 .doc 格式解析需要安装 LibreOffice。\n"
                "方案一：安装 LibreOffice (https://www.libreoffice.org/)\n"
                "方案二：将文件另存为 .docx 后重试"
            )
        tmp_dir = tempfile.mkdtemp()
        try:
            cmd = [lo_found, "--headless", "--convert-to", "docx", "--outdir", tmp_dir, str(path)]
            if password:
                cmd.insert(1, f"--infilter=MS Word 97:{password}")
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            if result.returncode != 0:
                raise DocumentLoaderError(f"LibreOffice 转换失败: {result.stderr}")
            converted = Path(tmp_dir) / f"{path.stem}.docx"
            if not converted.exists():
                raise DocumentLoaderError("转换后的文件未找到")
            return DocumentLoader._load_docx(converted, password)
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    @staticmethod
    def _load_pptx(path, password=None):
        from pptx import Presentation
        documents = []
        prs = Presentation(str(path))
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(cells))
            if not texts:
                continue
            documents.append({
                "page_content": "\n".join(texts),
                "metadata": {
                    "source": str(path), "filename": path.name,
                    "slide": slide_num, "total_slides": len(prs.slides),
                    "file_type": "pptx", "file_size": path.stat().st_size,
                },
            })
        return documents

    @staticmethod
    def _load_excel(path, password=None):
        ext = path.suffix.lower()
        all_sheets_text = []
        if ext == ".xlsx":
            try:
                import openpyxl
            except ImportError:
                raise DocumentLoaderError("需要安装 openpyxl")
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows_text = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        rows_text.append(" | ".join(str(cell) if cell is not None else "" for cell in row))
                if rows_text:
                    all_sheets_text.append(f"[工作表: {sheet_name}]\n" + "\n".join(rows_text))
            wb.close()
        elif ext == ".xls":
            try:
                import xlrd
            except ImportError:
                raise DocumentLoaderError("需要安装 xlrd")
            wb = xlrd.open_workbook(str(path))
            for sheet_name in wb.sheet_names():
                ws = wb.sheet_by_name(sheet_name)
                rows_text = []
                for row_idx in range(ws.nrows):
                    row = ws.row_values(row_idx)
                    rows_text.append(" | ".join(str(cell) if cell != "" else "" for cell in row))
                if rows_text:
                    all_sheets_text.append(f"[工作表: {sheet_name}]\n" + "\n".join(rows_text))
        if not all_sheets_text:
            return []
        return [{
            "page_content": "\n\n".join(all_sheets_text),
            "metadata": {
                "source": str(path), "filename": path.name,
                "file_type": ext.lstrip("."),
                "sheets": len(all_sheets_text),
                "file_size": path.stat().st_size,
            },
        }]

    @staticmethod
    def _load_image(path, password=None):
        from src.document_loader.ocr import ocr_image
        try:
            text = ocr_image(str(path))
        except ImportError as e:
            logger.warning(f"OCR 引擎未安装，图片跳过: {path.name} ({e})")
            return []
        except Exception as e:
            logger.warning(f"图片 OCR 识别失败，跳过: {path.name} ({e})")
            return []
        if not text.strip():
            logger.warning(f"图片 OCR 未识别到文字: {path.name}")
            return []
        return [{
            "page_content": text,
            "metadata": {
                "source": str(path), "filename": path.name,
                "file_type": path.suffix.lower().lstrip("."),
                "file_size": path.stat().st_size,
            },
        }]

    @staticmethod
    def _load_wps(path, password=None):
        ext = path.suffix.lower()
        try:
            subprocess.run(["libreoffice", "--version"], capture_output=True, timeout=10)
        except (FileNotFoundError, subprocess.TimeoutExpired):
            try:
                subprocess.run(["soffice", "--version"], capture_output=True, timeout=10)
            except (FileNotFoundError, subprocess.TimeoutExpired):
                raise DocumentLoaderError("WPS 格式需要 LibreOffice 转换支持")
        target_ext = ".docx" if ext == ".wps" else ".xlsx"
        tmp_dir = tempfile.mkdtemp()
        try:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", target_ext.lstrip("."),
                 "--outdir", tmp_dir, str(path)],
                capture_output=True, text=True, timeout=60,
            )
            if result.returncode != 0:
                raise DocumentLoaderError(f"LibreOffice 转换失败: {result.stderr}")
            converted = Path(tmp_dir) / f"{path.stem}{target_ext}"
            if not converted.exists():
                raise DocumentLoaderError("转换后的文件未找到")
            if target_ext == ".docx":
                return DocumentLoader._load_docx(converted)
            else:
                return DocumentLoader._load_excel(converted)
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    @staticmethod
    def _load_text(path, password=None):
        content = None
        used_encoding = "utf-8"
        raw_bytes = path.read_bytes()
        try:
            import chardet
            det = chardet.detect(raw_bytes)
            enc = det.get("encoding", "")
            conf = det.get("confidence", 0)
            if enc and conf > 0.5:
                try:
                    content = raw_bytes.decode(enc)
                    used_encoding = f"{enc} (chardet, {conf:.0%})"
                except (UnicodeDecodeError, LookupError):
                    pass
        except ImportError:
            pass
        if content is None:
            for enc in ["utf-8", "gbk", "gb2312", "utf-16", "big5", "shift_jis", "euc-kr"]:
                try:
                    content = raw_bytes.decode(enc)
                    used_encoding = enc
                    break
                except (UnicodeDecodeError, LookupError):
                    continue
        if content is None:
            content = raw_bytes.decode("latin-1")
            used_encoding = "latin-1 (保底)"
        return [{
            "page_content": content,
            "metadata": {
                "source": str(path), "filename": path.name,
                "file_type": path.suffix.lower().lstrip("."),
                "encoding": used_encoding, "file_size": path.stat().st_size,
            },
        }]

    @staticmethod
    def _load_json(path, password=None):
        try:
            data = json.loads(path.read_text(encoding="utf-8"))
        except UnicodeDecodeError:
            raw = path.read_bytes()
            import chardet
            det = chardet.detect(raw)
            enc = det.get("encoding", "utf-8") or "utf-8"
            data = json.loads(raw.decode(enc))
        except Exception as e:
            raise DocumentLoaderError(f"JSON 解析失败: {e}")
        return [{
            "page_content": json.dumps(data, ensure_ascii=False, indent=2),
            "metadata": {
                "source": str(path), "filename": path.name,
                "file_type": "json", "file_size": path.stat().st_size,
            },
        }]

    @property
    def stats(self):
        return {"loaded_files": self._load_count, "error_files": self._error_count}